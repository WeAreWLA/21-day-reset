"""
Build the WLA Founding Members launch presentation.
16:9 widescreen, brand-matched to the /founding-members landing page:
  - cream paper background, navy serif headlines, coral italic accents
  - Libre Baskerville (serif) + Alegreya Sans (body)
  - callout slides on navy

Output: wla-summer-challenge-launch.pptx (import to Google Slides via
File -> Import slides).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------- Brand ----------
NAVY        = RGBColor(0x0E, 0x27, 0x46)
NAVY_DEEP   = RGBColor(0x09, 0x1A, 0x30)
CORAL       = RGBColor(0xD9, 0x77, 0x57)
CORAL_SOFT  = RGBColor(0xE5, 0x91, 0x7A)
CREAM       = RGBColor(0xF9, 0xF7, 0xF4)
CREAM_2     = RGBColor(0xE9, 0xDF, 0xD3)
PAPER       = RGBColor(0xFD, 0xFB, 0xF8)
INK         = RGBColor(0x1A, 0x24, 0x38)
MUTED       = RGBColor(0x5A, 0x64, 0x78)
BODY        = RGBColor(0x33, 0x33, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_WARM  = RGBColor(0xF4, 0xEB, 0xDC)

SERIF = 'Libre Baskerville'
SANS  = 'Alegreya Sans'

# ---------- Presentation ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, left, top, width, height, text,
             font=SANS, size=18, color=BODY, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.2,
             letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if letter_spacing is not None:
        # OOXML run property w:spc isn't directly exposed; use char spacing in 100ths of a point
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(letter_spacing))
    return tb, tf, p, r

def add_runs(slide, left, top, width, height, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.2):
    """runs: list of dicts {text, font, size, color, bold, italic, letter_spacing}"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    for i, run in enumerate(runs):
        r = p.add_run()
        r.text = run['text']
        r.font.name  = run.get('font', SANS)
        r.font.size  = Pt(run.get('size', 18))
        r.font.bold  = run.get('bold', False)
        r.font.italic = run.get('italic', False)
        r.font.color.rgb = run.get('color', BODY)
        if 'letter_spacing' in run:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(run['letter_spacing']))
    return tb

def add_paragraphs(slide, left, top, width, height, paragraphs,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of {runs:[...], align?, line?, space_after?}"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align)
        p.line_spacing = para.get('line', 1.25)
        if 'space_after' in para:
            p.space_after = Pt(para['space_after'])
        for run in para['runs']:
            r = p.add_run()
            r.text = run['text']
            r.font.name  = run.get('font', SANS)
            r.font.size  = Pt(run.get('size', 18))
            r.font.bold  = run.get('bold', False)
            r.font.italic = run.get('italic', False)
            r.font.color.rgb = run.get('color', BODY)
    return tb

def add_rect(slide, left, top, width, height, color, line_color=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width is not None:
            s.line.width = line_width
    s.shadow.inherit = False
    return s

def add_round_rect(slide, left, top, width, height, color, radius_pct=0.06, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # adjust the corner radius via the adj1 value (0 to ~0.5)
    s.adjustments[0] = radius_pct
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
    s.shadow.inherit = False
    return s

def add_line(slide, x1, y1, x2, y2, color, weight_pt=1.5):
    from pptx.shapes.connector import Connector
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line

def add_pill(slide, left, top, width, height, text, fill=CORAL, color=WHITE, size=11,
             font=SANS, bold=True, letter_spacing=180):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    rPr.set('spc', str(letter_spacing))
    return pill

def add_dot(slide, cx, cy, d, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(int(cx - d/2)), Emu(int(cy - d/2)),
        Emu(int(d)), Emu(int(d)))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return dot

# ---------- Reusable slide patterns ----------

def slide_eyebrow_headline(slide, eyebrow, title_runs, body_paras=None,
                            top=Inches(0.85), left=Inches(0.9), width=Inches(11.5)):
    """Eyebrow text + headline (rich runs) + optional body paragraphs. Top-aligned by default."""
    if eyebrow:
        add_runs(slide, left, top, width, Inches(0.55), [{
            'text': eyebrow.upper(), 'font': SANS, 'size': 18, 'color': CORAL,
            'bold': True, 'letter_spacing': 220
        }])
    add_paragraphs(slide, left, top + Inches(0.85), width, Inches(2.5),
        [{'runs': title_runs, 'line': 1.12, 'align': PP_ALIGN.LEFT}])
    if body_paras:
        add_paragraphs(slide, left, top + Inches(2.9), width, Inches(3.5),
            body_paras)

def coral_bullet(slide, left, top, text, size=18, font=SERIF, color=NAVY, italic=False, bold=False):
    """Draw a coral dot + a bit of text. Returns the y of the next line position."""
    dot_d = Emu(Inches(0.13))
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Emu(int(Pt(size).emu * 0.45)), dot_d, dot_d)
    dot.fill.solid(); dot.fill.fore_color.rgb = CORAL
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_runs(slide, left + Inches(0.34), top, Inches(11.0), Inches(0.55),
        [{'text': text, 'font': font, 'size': size, 'color': color, 'italic': italic, 'bold': bold}])

def add_footer(slide, dark=False):
    color = WHITE if dark else MUTED
    add_runs(slide, Inches(0.6), Inches(6.95), Inches(8), Inches(0.4),
        [{'text': 'The Weight Loss Academy', 'font': SERIF, 'size': 13,
          'color': color, 'italic': True}])
    add_runs(slide, Inches(7.5), Inches(6.95), Inches(5.3), Inches(0.4),
        [{'text': 'WLA SUMMER CHALLENGE  ·  FOUNDING MEMBERS', 'font': SANS,
          'size': 12, 'color': color, 'bold': True, 'letter_spacing': 200}],
        align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDES
# ============================================================

# ---------------- 1. COVER ----------------
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    # Soft radial-ish corner blocks via large rounded rects (approximation)
    big = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3), Inches(7), Inches(7))
    big.fill.solid(); big.fill.fore_color.rgb = CREAM_2; big.line.fill.background()
    big.shadow.inherit = False
    big2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(7), Inches(7))
    big2.fill.solid(); big2.fill.fore_color.rgb = CREAM_2; big2.line.fill.background()
    big2.shadow.inherit = False

    add_pill(s, Inches(3.4), Inches(0.9), Inches(6.5), Inches(0.55),
             '✦  SOMETHING BIG IS HAPPENING INSIDE WLA  ✦',
             fill=NAVY, color=CREAM_WARM, size=15, letter_spacing=260)

    add_paragraphs(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.5), [
        {'runs': [
            {'text': 'The Next Phase ', 'font': SERIF, 'size': 64, 'color': NAVY, 'bold': True},
            {'text': 'of', 'font': SERIF, 'size': 64, 'color': NAVY, 'italic': True},
         ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
        {'runs': [
            {'text': 'The Weight Loss Academy', 'font': SERIF, 'size': 64, 'color': CORAL, 'italic': True},
         ], 'align': PP_ALIGN.CENTER, 'line': 1.1, 'space_after': 24},
    ])

    add_runs(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.7),
        [{'text': 'Lose the weight. Live the lifestyle.',
          'font': SERIF, 'size': 30, 'color': NAVY, 'italic': True}],
        align=PP_ALIGN.CENTER)

    add_runs(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.5),
        [{'text': 'A live announcement with Anna Wallace  ·  Founder of WLA',
          'font': SANS, 'size': 18, 'color': MUTED, 'bold': True, 'letter_spacing': 200}],
        align=PP_ALIGN.CENTER)

# ---------------- 2. OPENING ----------------
def slide_opening():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'A quick word before we begin',
        [{'text': 'Today I want to share something ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'really exciting', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
         {'text': ' with you.', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}],
        body_paras=[
            {'runs': [{'text': "Over the last few months, behind the scenes, we've been building something that I genuinely believe is going to change the future of WLA.",
                       'font': SANS, 'size': 19, 'color': BODY}], 'line': 1.5, 'space_after': 10},
            {'runs': [{'text': 'And honestly… I think this is going to change the game for so many women.',
                       'font': SERIF, 'size': 20, 'color': NAVY, 'italic': True}], 'line': 1.5, 'space_after': 10},
            {'runs': [{'text': "But before I reveal it… I want to talk about something important.",
                       'font': SANS, 'size': 19, 'color': BODY}], 'line': 1.5},
        ])
    add_footer(s)

# ---------------- 3. 11 YEARS ----------------
def slide_11_years():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'Over the last 11 years',
        [{'text': 'WLA has helped over ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': '50,000 women', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
         {'text': '…', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}],
        top=Inches(0.85))

    # two columns of bullets
    bullets = [
        'lose weight',
        'reduce cravings',
        'feel more confident',
        'build healthier habits',
        'improve their health',
        'put type 2 diabetes into remission',
        'and finally feel back in control around food again',
    ]
    y = Inches(3.8)
    for i, b in enumerate(bullets):
        col = i % 2
        row = i // 2
        x = Inches(0.95) if col == 0 else Inches(7.0)
        ty = y + Inches(0.62 * row)
        coral_bullet(s, x, ty, b.capitalize(), size=18, color=NAVY)
    add_footer(s)

# ---------------- 4. THE HARDEST PART ----------------
def slide_hardest_part():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'But one thing became very clear',
        [{'text': 'The hardest part of weight loss is ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
         {'text': 'rarely starting.', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True}],
        body_paras=[
            {'runs': [{'text': 'The hardest part… is maintaining the results once life gets busy again.',
                       'font': SERIF, 'size': 22, 'color': NAVY, 'italic': True}], 'line': 1.5, 'space_after': 14},
            {'runs': [{'text': 'Most women do not struggle to LOSE weight. They struggle to KEEP it off once accountability disappears, routines slowly fade, old habits creep back in and life takes over again.',
                       'font': SANS, 'size': 18, 'color': BODY}], 'line': 1.55},
        ],
        top=Inches(0.85))
    add_footer(s)

# ---------------- 5. THE CYCLE (CALLOUT NAVY) ----------------
def slide_cycle():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_pill(s, Inches(5.2), Inches(1.4), Inches(2.9), Inches(0.4),
             'THE CYCLE', fill=CORAL, color=WHITE, size=15, letter_spacing=300)
    add_paragraphs(s, Inches(1.2), Inches(2.4), Inches(10.9), Inches(3.5), [
        {'runs': [
            {'text': '“I’ll start again ', 'font': SERIF, 'size': 64, 'color': WHITE, 'italic': True},
            {'text': 'Monday', 'font': SERIF, 'size': 64, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
            {'text': '.”', 'font': SERIF, 'size': 64, 'color': WHITE, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_runs(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.6),
        [{'text': "And THAT is the cycle we want to help women finally break for good.",
          'font': SANS, 'size': 18, 'color': CREAM_WARM}],
        align=PP_ALIGN.CENTER)
    add_footer(s, dark=True)

# ---------------- 6. WOMEN WHO KEEP IT OFF ----------------
def slide_keep_it_off():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'The women who keep the weight off',
        [{'text': 'Usually do ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'one thing', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
         {'text': ' differently.', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}],
        body_paras=[
            {'runs': [{'text': 'They continue the lifestyle.',
                       'font': SERIF, 'size': 24, 'color': NAVY, 'italic': True}], 'line': 1.5, 'space_after': 16},
            {'runs': [{'text': 'They stay connected to the support, accountability, routines, coaching and habits that helped create their results in the first place.',
                       'font': SANS, 'size': 18, 'color': BODY}], 'line': 1.55},
        ])
    add_footer(s)

# ---------------- 7. THE WLA PILLARS ----------------
def slide_pillars():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_runs(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4),
        [{'text': 'WHAT THE LIFESTYLE LOOKS LIKE', 'font': SANS, 'size': 18,
          'color': CORAL, 'bold': True, 'letter_spacing': 220}])
    add_paragraphs(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(2.0),
        [{'runs': [
            {'text': 'They build a ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
            {'text': 'lifestyle', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True},
            {'text': ' that supports their goals.', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
        ], 'line': 1.1}])

    pillars = ['Structure', 'Accountability', 'Routines', 'Real support', 'Flexibility', 'Consistency']
    col_w = Inches(3.9)
    row_h = Inches(1.1)
    start_x = Inches(0.9); start_y = Inches(3.6)
    gap = Inches(0.15)
    for i, p in enumerate(pillars):
        col = i % 3; row = i // 3
        x = start_x + (col_w + gap) * col
        y = start_y + (row_h + gap) * row
        card = add_round_rect(s, x, y, col_w, row_h, PAPER, radius_pct=0.18)
        card.line.color.rgb = CREAM_2
        # coral left bar
        bar = add_rect(s, x, y, Emu(int(Inches(0.07).emu)), row_h, CORAL)
        # text
        add_runs(s, x + Inches(0.32), y, col_w - Inches(0.32), row_h,
                 [{'text': p, 'font': SERIF, 'size': 22, 'color': NAVY, 'bold': True}],
                 anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s)

# ---------------- 8. THE EVOLUTION ----------------
def slide_evolution():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'The evolution of WLA',
        [{'text': 'The resets and challenges were ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
         {'text': 'never the finish line.', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True}],
        body_paras=[
            {'runs': [{'text': 'They were designed to create momentum, help women reset, reduce cravings, rebuild consistency and create results quickly. And they work incredibly well for that.',
                       'font': SANS, 'size': 18, 'color': BODY}], 'line': 1.55, 'space_after': 14},
            {'runs': [{'text': 'But the women who achieve the best long-term results are usually the women who continue the habits and routines long after the challenge ends.',
                       'font': SERIF, 'size': 20, 'color': NAVY, 'italic': True}], 'line': 1.5},
        ],
        top=Inches(0.85))
    add_footer(s)

# ---------------- 9. THE SHIFT (CALLOUT) ----------------
def slide_shift():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_pill(s, Inches(5.7), Inches(1.5), Inches(1.9), Inches(0.4),
             'THE SHIFT', fill=CORAL, color=WHITE, size=15, letter_spacing=300)
    add_paragraphs(s, Inches(1.2), Inches(2.5), Inches(10.9), Inches(3.5), [
        {'runs': [
            {'text': 'Not just helping women ', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True},
            {'text': 'lose', 'font': SERIF, 'size': 46, 'color': CORAL_SOFT, 'italic': True},
            {'text': ' the weight…', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
        {'runs': [
            {'text': 'But helping women ', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True},
            {'text': 'keep', 'font': SERIF, 'size': 46, 'color': CORAL_SOFT, 'italic': True},
            {'text': ' it off too.', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_footer(s, dark=True)

# ---------------- 10. QUOTE CALLOUT ----------------
def slide_quote_lifestyle():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_runs(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.4),
        [{'text': 'REMEMBER THIS', 'font': SANS, 'size': 18,
          'color': CORAL, 'bold': True, 'letter_spacing': 240}])
    # Big serif quote
    add_paragraphs(s, Inches(1.2), Inches(2.1), Inches(10.9), Inches(4), [
        {'runs': [
            {'text': '“The women who keep the weight off', 'font': SERIF, 'size': 48, 'color': NAVY, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
        {'runs': [
            {'text': 'usually ', 'font': SERIF, 'size': 48, 'color': NAVY, 'italic': True},
            {'text': 'continue the lifestyle', 'font': SERIF, 'size': 48, 'color': CORAL, 'italic': True, 'bold': True},
            {'text': '.”', 'font': SERIF, 'size': 48, 'color': NAVY, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_runs(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
        [{'text': '— Anna Wallace, Founder of WLA',
          'font': SANS, 'size': 16, 'color': MUTED, 'bold': True, 'letter_spacing': 200}],
        align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 11. I LIVE IT TOO ----------------
def slide_i_live_it():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM_2)
    add_runs(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.4),
        [{'text': 'AND HONESTLY…', 'font': SANS, 'size': 18,
          'color': CORAL, 'bold': True, 'letter_spacing': 240}])
    add_paragraphs(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4), [
        {'runs': [
            {'text': 'I still live the ', 'font': SERIF, 'size': 56, 'color': NAVY, 'bold': True},
            {'text': 'WLA way', 'font': SERIF, 'size': 56, 'color': CORAL, 'italic': True},
        ], 'line': 1.1},
        {'runs': [
            {'text': 'every single day too.', 'font': SERIF, 'size': 56, 'color': NAVY, 'bold': True},
        ], 'line': 1.1},
    ])
    add_footer(s)

# ---------------- 12. NOT PERFECTLY ----------------
def slide_not_perfectly():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'Not perfectly',
        [{'text': 'But ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'consistently enough', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
         {'text': ' to:', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}],
        top=Inches(0.85))
    bullets = [
        'maintain my results',
        'feel good in myself',
        'stay in control around food',
        'and build routines that fit real life',
    ]
    y = Inches(3.7)
    for i, b in enumerate(bullets):
        coral_bullet(s, Inches(1.4), y + Inches(0.62 * i), b.capitalize(), size=20)
    add_footer(s)

# ---------------- 13. I STILL ----------------
def slide_i_still():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'What that looks like for me',
        [{'text': 'I still…', 'font': SERIF, 'size': 56, 'color': NAVY, 'bold': True}],
        top=Inches(0.85))
    bullets = [
        'plan my meals',
        'focus on simple nutrition',
        'enjoy indulgences',
        'move my body',
        'prioritise routines',
        'and use the exact same principles we teach inside WLA',
    ]
    for i, b in enumerate(bullets):
        col = i % 2; row = i // 2
        x = Inches(1.0) if col == 0 else Inches(7.0)
        y = Inches(3.4) + Inches(0.65 * row)
        coral_bullet(s, x, y, b.capitalize(), size=18, color=NAVY)
    add_footer(s)

# ---------------- 14. NEVER A SHORT-TERM DIET ----------------
def slide_never_diet():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'WLA was never meant to be',
        [{'text': 'Another ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'short-term diet.', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True}],
        body_paras=[
            {'runs': [{'text': 'It was designed to help women build a lifestyle they can actually LIVE with long term.',
                       'font': SERIF, 'size': 22, 'color': NAVY, 'italic': True}], 'line': 1.5},
        ],
        top=Inches(0.85))
    add_footer(s)

# ---------------- 15. NOT TAUGHT - LIVED (CALLOUT) ----------------
def slide_taught_lived():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_paragraphs(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'WLA isn’t something I just ', 'font': SERIF, 'size': 50, 'color': WHITE, 'bold': True},
            {'text': 'teach', 'font': SERIF, 'size': 50, 'color': CORAL_SOFT, 'italic': True},
            {'text': '.', 'font': SERIF, 'size': 50, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
        {'runs': [
            {'text': 'It’s something I ', 'font': SERIF, 'size': 50, 'color': WHITE, 'bold': True},
            {'text': 'live', 'font': SERIF, 'size': 50, 'color': CORAL_SOFT, 'italic': True},
            {'text': '.', 'font': SERIF, 'size': 50, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_footer(s, dark=True)

# ---------------- 16. PASSIONATE ----------------
def slide_passionate():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_paragraphs(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'And honestly… that is why I’m so ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
            {'text': 'passionate', 'font': SERIF, 'size': 38, 'color': CORAL, 'italic': True},
            {'text': ' about this', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
        {'runs': [
            {'text': 'next phase of WLA.', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_footer(s)

# ---------------- 17. WELCOME TO NEXT PHASE ----------------
def slide_welcome():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(4.6), Inches(1.2), Inches(4.1), Inches(0.4),
             'WELCOME TO THE NEXT PHASE OF WLA',
             fill=CREAM_2, color=CORAL, size=15, letter_spacing=300)
    add_paragraphs(s, Inches(0.5), Inches(2.4), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'The Weight Loss', 'font': SERIF, 'size': 60, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
        {'runs': [
            {'text': '& Lifestyle ', 'font': SERIF, 'size': 60, 'color': NAVY, 'bold': True},
            {'text': 'Academy', 'font': SERIF, 'size': 60, 'color': CORAL, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
    ])
    add_runs(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6),
        [{'text': 'Lose the weight. Live the lifestyle.',
          'font': SERIF, 'size': 24, 'color': NAVY, 'italic': True}],
        align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 18. WHAT IT'S ABOUT ----------------
def slide_about():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'This next chapter is about more',
        [{'text': 'About helping you ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
         {'text': 'keep going', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True},
         {'text': '.', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True}],
        top=Inches(0.85))
    bullets = [
        'maintain momentum',
        'stay consistent',
        'build routines that LAST',
        'and finally stop starting over',
    ]
    y = Inches(3.4)
    for i, b in enumerate(bullets):
        coral_bullet(s, Inches(2.0), y + Inches(0.65 * i), b.capitalize(), size=22)
    add_footer(s)

# ---------------- 19. SOMETHING HUGE BEHIND THE SCENES ----------------
def slide_huge():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, "And that’s why we’ve built something huge…",
        [{'text': 'Designed to help you ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
         {'text': 'continue the WLA lifestyle.', 'font': SERIF, 'size': 38, 'color': CORAL, 'italic': True}],
        top=Inches(0.85))
    bullets = [
        'stay accountable',
        'stay focused',
        'continue the WLA lifestyle',
        'and maintain results long after the initial weight loss phase',
    ]
    y = Inches(3.8)
    for i, b in enumerate(bullets):
        coral_bullet(s, Inches(2.0), y + Inches(0.6 * i), b.capitalize(), size=18)
    add_footer(s)

# ---------------- 20. REVEAL: THE APP ----------------
def slide_reveal_app():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_pill(s, Inches(4.6), Inches(1.2), Inches(4.1), Inches(0.4),
             '✨  INTRODUCING  ✨', fill=CORAL, color=WHITE, size=15, letter_spacing=320)
    add_paragraphs(s, Inches(0.5), Inches(2.5), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'The Brand New ', 'font': SERIF, 'size': 70, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
        {'runs': [
            {'text': 'WLA App', 'font': SERIF, 'size': 90, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.05},
    ])
    add_runs(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
        [{'text': 'Support in your pocket. Every single day.',
          'font': SERIF, 'size': 22, 'color': CREAM_WARM, 'italic': True}],
        align=PP_ALIGN.CENTER)
    add_footer(s, dark=True)

# ---------------- 21. NOT ANOTHER CALORIE APP ----------------
def slide_not_calorie():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'But first…',
        [{'text': 'This is ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'not', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
         {'text': ' another calorie-counting app.', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}],
        top=Inches(0.85))
    add_runs(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.5),
        [{'text': 'It’s been built to help you:', 'font': SERIF, 'size': 22,
          'color': NAVY, 'italic': True}])
    bullets = [
        'hold onto your results',
        'continue to get consistent results',
        'stay accountable & focused',
        'continue the lifestyle',
        'maintain momentum',
        'and stop falling off track every time life gets busy',
    ]
    for i, b in enumerate(bullets):
        col = i % 2; row = i // 2
        x = Inches(1.0) if col == 0 else Inches(7.0)
        y = Inches(4.4) + Inches(0.55 * row)
        coral_bullet(s, x, y, b.capitalize(), size=16, color=NAVY)
    add_footer(s)

# ---------------- 22. WHAT'S INSIDE (OVERVIEW) ----------------
def slide_inside_overview():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'Inside the app you get access to',
        [{'text': 'Everything ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
         {'text': 'in one place.', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True}],
        top=Inches(0.85))
    items = [
        'Accountability & progress tracking',
        'Meal planning & shopping lists',
        'Recipes & quick meals',
        'Coaching & support',
        'Lifestyle tracking',
        'Community support',
    ]
    for i, b in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.9) if col == 0 else Inches(7.0)
        y = Inches(3.4) + Inches(0.7 * row)
        coral_bullet(s, x, y, b, size=20, color=NAVY)
    add_footer(s)

# ---------------- 23. PHASE 1 INTRO ----------------
def slide_phase1_intro():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(4.6), Inches(1.4), Inches(4.1), Inches(0.4),
             '🚀  AVAILABLE INSIDE PHASE 1 AT LAUNCH',
             fill=NAVY, color=CREAM_WARM, size=15, letter_spacing=260)
    add_paragraphs(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'Here’s what Founding Members get ', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
            {'text': 'immediately', 'font': SERIF, 'size': 42, 'color': CORAL, 'italic': True},
            {'text': '.', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_runs(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6),
        [{'text': 'From the moment you join, you unlock the very first version of the new WLA experience.',
          'font': SANS, 'size': 18, 'color': BODY}],
        align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 24–33. PHASE 1 FEATURE SLIDES ----------------
def feature_slide(num, total, title, lead, bullets, footer_quote=None):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    # Top-left feature counter
    add_runs(s, Inches(0.9), Inches(0.85), Inches(6), Inches(0.5),
        [{'text': f'WLA APP FEATURE  ·  {num:02d} / {total:02d}',
          'font': SANS, 'size': 20, 'color': CORAL, 'bold': True, 'letter_spacing': 200}])
    # Star + title
    add_runs(s, Inches(0.9), Inches(1.55), Inches(0.6), Inches(0.6),
        [{'text': '✦', 'font': SERIF, 'size': 32, 'color': CORAL}])
    add_paragraphs(s, Inches(1.55), Inches(1.5), Inches(11.0), Inches(1.3), [
        {'runs': [
            {'text': title, 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
        ], 'line': 1.15}
    ])
    if lead:
        add_runs(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.6),
            [{'text': lead, 'font': SERIF, 'size': 22, 'color': NAVY, 'italic': True}])
    # Two-column bullets
    start_y = Inches(3.95)
    half = (len(bullets) + 1) // 2
    for i, b in enumerate(bullets):
        col = 0 if i < half else 1
        row = i if col == 0 else (i - half)
        x = Inches(0.9) if col == 0 else Inches(7.0)
        y = start_y + Inches(0.5 * row)
        coral_bullet(s, x, y, b, size=18, color=NAVY)
    if footer_quote:
        add_runs(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.4),
            [{'text': footer_quote, 'font': SERIF, 'size': 17, 'color': CORAL, 'italic': True}],
            align=PP_ALIGN.LEFT)
    add_footer(s)

PHASE1_FEATURES = [
    {
        'title': 'Daily Accountability + Progress Tracking',
        'lead':  'Stay focused without feeling overwhelmed.',
        'bullets': ['Meals', 'Habits', 'Routines', 'Progress', 'Consistency'],
        'footer': 'No more “I’ve completely fallen off track.”',
    },
    {
        'title': 'Weight Loss Dashboard + Goal Tracking',
        'lead':  'Stay motivated and watch the momentum build.',
        'bullets': ['Weight tracking', 'Visual progress graphs', 'Milestone tracking', 'Monthly check-ins', 'Progress dashboards'],
        'footer': 'See how far you’ve come and keep momentum high.',
    },
    {
        'title': 'Smart Weekly Meal Planner',
        'lead':  'Plan your week in minutes.',
        'bullets': ['Drag meals into breakfast, lunch & dinner', 'Know exactly what you’re eating', 'Know exactly what to buy', 'Stay organised during busy weeks'],
        'footer': 'Consistency becomes so much easier with a plan in place.',
    },
    {
        'title': 'Automatic Shopping Lists',
        'lead':  'Your shopping list, written for you.',
        'bullets': ['Generated from your meal planner', 'Less stress, less overthinking', 'More staying on track'],
        'footer': 'One tap. Sorted.',
    },
    {
        'title': 'Weekly Meal Guides + Fresh Recipes Every Friday',
        'lead':  'Never feel stuck for ideas again.',
        'bullets': ['Fresh recipes', 'Meal inspiration', 'Simple nutrition guidance', 'Quick meals & no-cook options', 'Realistic ideas for busy lifestyles'],
        'footer': 'New content added directly into the app every single week.',
    },
    {
        'title': 'Extensive WLA Recipe Library',
        'lead':  '100s of delicious recipes built for real life.',
        'bullets': ['Weight loss', 'Busy women', 'Family life', 'Results without restriction'],
        'footer': 'No bland diet food. No miserable meal plans. No living on salads.',
    },
    {
        'title': 'Lifestyle Tracking For Better Results',
        'lead':  'Because weight loss is about more than food alone.',
        'bullets': ['Sleep', 'Stress', 'Movement', 'Water intake', 'Alcohol', 'Caffeine', 'Mood', 'Routines'],
        'footer': 'Identify the habits helping (or hurting) your results.',
    },
    {
        'title': 'Habit Tracking System',
        'lead':  'Build habits that actually LAST beyond motivation.',
        'bullets': ['Routines', 'Structure', 'Consistency', 'Continuing the lifestyle long after the challenge ends'],
        'footer': 'Long-term results are built habit by habit.',
    },
    {
        'title': 'Weekly Coaching Sessions With Anna',
        'lead':  'Coaching, accountability and momentum every week.',
        'bullets': ['Coaching', 'Accountability', 'Support', 'Guidance', 'Motivation every week'],
        'footer': 'Always knowing what to focus on next.',
    },
    {
        'title': 'Female-Only WLA Community',
        'lead':  'Stay connected with women on the same journey.',
        'bullets': ['Share wins', 'Celebrate progress', 'Stay accountable', 'Stop trying to do this alone'],
        'footer': 'Like having support in your pocket every single day.',
    },
]

# ---------------- 34. PLUS, ONLY THE BEGINNING ----------------
def slide_plus_beginning():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM_2)
    add_pill(s, Inches(4.2), Inches(1.2), Inches(4.9), Inches(0.4),
             '✨  PLUS… THIS IS ONLY THE BEGINNING',
             fill=NAVY, color=CREAM_WARM, size=15, letter_spacing=280)
    add_paragraphs(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'Founding Members get ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
            {'text': 'everything coming next', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True},
            {'text': ',', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
        {'runs': [
            {'text': 'too.', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_runs(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6),
        [{'text': 'Future app updates and new features as WLA continues expanding, before our public Phase 2 launch in September.',
          'font': SANS, 'size': 17, 'color': BODY}], align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 35. COMING IN PHASE 2 INTRO ----------------
def slide_phase2_intro():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(4.7), Inches(1.4), Inches(3.9), Inches(0.4),
             '🚀  COMING SOON · PHASE 2',
             fill=CORAL, color=WHITE, size=15, letter_spacing=280)
    add_paragraphs(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': "What’s coming next is ", 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
            {'text': 'very', 'font': SERIF, 'size': 42, 'color': CORAL, 'italic': True},
            {'text': ' exciting.', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_footer(s)

# ---------------- 36–39. PHASE 2 FEATURES ----------------
def phase2_slide(num, title, lead, bullets):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_runs(s, Inches(0.9), Inches(0.85), Inches(6), Inches(0.5),
        [{'text': f'COMING IN PHASE 2  ·  {num:02d} / 04',
          'font': SANS, 'size': 20, 'color': CORAL, 'bold': True, 'letter_spacing': 200}])
    add_runs(s, Inches(0.9), Inches(1.55), Inches(0.6), Inches(0.6),
        [{'text': '🚀', 'font': SERIF, 'size': 28, 'color': CORAL}])
    add_paragraphs(s, Inches(1.55), Inches(1.5), Inches(11.0), Inches(1.3),
        [{'runs': [{'text': title, 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True}], 'line': 1.15}])
    if lead:
        add_runs(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.6),
            [{'text': lead, 'font': SERIF, 'size': 22, 'color': NAVY, 'italic': True}])
    y = Inches(3.95)
    for i, b in enumerate(bullets):
        coral_bullet(s, Inches(1.2), y + Inches(0.6 * i), b, size=18, color=NAVY)
    add_footer(s)

PHASE2_FEATURES = [
    {'title': 'Personalised Protein + Fibre Targets',
     'lead':  'Make consistency easier than ever.',
     'bullets': ['Reduce cravings', 'Stay fuller for longer', 'Simplify nutrition', 'Make consistency easier than ever before']},
    {'title': 'The Updated WLA Formula',
     'lead':  'The next evolution of the WLA approach.',
     'bullets': ['Lose weight', 'Maintain results', 'Simplify nutrition', 'Build a lifestyle that actually LASTS']},
    {'title': 'Meal Photo Nutrition Analysis',
     'lead':  'Snap a photo. Get guidance.',
     'bullets': ['Quicker', 'Easier', 'More realistic for busy lives']},
    {'title': 'On-Demand Support Inside The App',
     'lead':  'Support exactly when you need it most.',
     'bullets': ['Feeling off track', 'Struggling with cravings', 'Needing support during a difficult week']},
]

# ---------------- 40. THIS IS THE NEXT EVOLUTION (CALLOUT) ----------------
def slide_next_evolution():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_pill(s, Inches(5.0), Inches(1.4), Inches(3.3), Inches(0.4),
             'THE NEXT EVOLUTION', fill=CORAL, color=WHITE, size=15, letter_spacing=280)
    add_paragraphs(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'Not just helping women ', 'font': SERIF, 'size': 42, 'color': WHITE, 'bold': True},
            {'text': 'lose weight', 'font': SERIF, 'size': 42, 'color': CORAL_SOFT, 'italic': True},
            {'text': '…', 'font': SERIF, 'size': 42, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.18},
        {'runs': [
            {'text': 'But helping them ', 'font': SERIF, 'size': 42, 'color': WHITE, 'bold': True},
            {'text': 'keep it off too.', 'font': SERIF, 'size': 42, 'color': CORAL_SOFT, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.18},
    ])
    add_footer(s, dark=True)

# ---------------- 41. FOUNDING MEMBERS OPEN ----------------
def slide_founding_open():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(3.9), Inches(1.0), Inches(5.5), Inches(0.45),
             '🚨  FOUNDING MEMBERS ARE NOW OPENING  🚨',
             fill=CORAL, color=WHITE, size=16, letter_spacing=260)
    add_paragraphs(s, Inches(0.5), Inches(2.4), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'Only ', 'font': SERIF, 'size': 60, 'color': NAVY, 'bold': True},
            {'text': '1,000 places', 'font': SERIF, 'size': 60, 'color': CORAL, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
        {'runs': [
            {'text': 'before public launch.', 'font': SERIF, 'size': 60, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.1},
    ])
    add_runs(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.6),
        [{'text': 'Because you’re already part of WLA, YOU get first access before anyone else.',
          'font': SERIF, 'size': 19, 'color': NAVY, 'italic': True}], align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 42. ONCE THEY'RE GONE (CALLOUT) ----------------
def slide_gone_forever():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_paragraphs(s, Inches(0.5), Inches(2.6), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'Once these places are gone…', 'font': SERIF, 'size': 44, 'color': WHITE, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
        {'runs': [
            {'text': 'they are gone ', 'font': SERIF, 'size': 56, 'color': WHITE, 'bold': True},
            {'text': 'forever', 'font': SERIF, 'size': 56, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
            {'text': '.', 'font': SERIF, 'size': 56, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_footer(s, dark=True)

# ---------------- 43. ONLY OPPORTUNITY ----------------
def slide_only_opportunity():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'Once Phase 1 closes',
        [{'text': 'Founding Member access ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
         {'text': 'disappears.', 'font': SERIF, 'size': 38, 'color': CORAL, 'italic': True}],
        top=Inches(0.85))
    bullets = [
        'Founding Member access disappears',
        'Pricing increases',
        'This exact offer will never be available again',
    ]
    for i, b in enumerate(bullets):
        coral_bullet(s, Inches(1.5), Inches(3.2) + Inches(0.65 * i), b, size=20)
    # Founding members will receive label
    add_runs(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.4),
        [{'text': 'FOUNDING MEMBERS WILL RECEIVE', 'font': SANS, 'size': 17,
          'color': CORAL, 'bold': True, 'letter_spacing': 220}])
    receives = [
        'The LOWEST price WLA will ever be offered at',
        'Early access before public launch',
        'FREE access to the new WLA app',
        'Future updates and new features',
        'A place inside the future of WLA from the very beginning',
    ]
    for i, b in enumerate(receives):
        col = i % 2; row = i // 2
        x = Inches(0.9) if col == 0 else Inches(7.0)
        y = Inches(5.95) + Inches(0.45 * row)
        # check icon
        chk = slide_check_icon(s, x, y + Inches(0.08), size=Inches(0.24))
        add_runs(s, x + Inches(0.40), y, Inches(5.8), Inches(0.5),
            [{'text': b, 'font': SANS, 'size': 17, 'color': NAVY, 'bold': True}])

def slide_check_icon(s, x, y, size=Inches(0.18), fill=CORAL):
    box = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '✓'
    r.font.name = SANS; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = WHITE
    return box

# ---------------- 44. WHAT FOUNDING MEMBERS RECEIVE ----------------
def slide_what_receive():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'So here’s what Founding Members receive',
        [{'text': 'Everything you need to ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
         {'text': 'actually keep going.', 'font': SERIF, 'size': 38, 'color': CORAL, 'italic': True}],
        top=Inches(0.85))
    items = [
        ('The 12 Week WLA Summer Challenge', 'Designed for women 45+ to lose, feel and finally maintain.'),
        ('FREE 12 month access to the new WLA App', 'Coaching, planning, recipes, tracking — all in one place.'),
        ('Access to ALL WLA Challenges for 12 months', 'Every reset, challenge and accountability event.'),
        ('Weekly Coaching With Anna', 'Support, guidance and accountability every single week.'),
        ('Female-Only Community Support', 'Daily accountability and women on the same journey.'),
    ]
    y_start = Inches(3.1)
    for i, (t, d) in enumerate(items):
        y = y_start + Inches(0.72 * i)
        chk = slide_check_icon(s, Inches(0.95), y + Inches(0.10), size=Inches(0.32))
        add_runs(s, Inches(1.5), y, Inches(11), Inches(0.4),
            [{'text': t, 'font': SERIF, 'size': 19, 'color': NAVY, 'bold': True}])
        add_runs(s, Inches(1.5), y + Inches(0.36), Inches(11), Inches(0.4),
            [{'text': d, 'font': SANS, 'size': 16, 'color': BODY}])
    add_footer(s)

# ---------------- 45. PRICING SLIDE ----------------
def slide_pricing():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(4.5), Inches(0.8), Inches(4.3), Inches(0.4),
             'FOUNDING MEMBER PRICING', fill=CREAM_2, color=CORAL, size=15, letter_spacing=300)
    # The price card
    card = add_round_rect(s, Inches(2.5), Inches(1.6), Inches(8.3), Inches(4.6),
                          PAPER, radius_pct=0.06)
    card.line.color.rgb = CREAM_2
    # All for only
    add_runs(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.5),
        [{'text': 'ALL FOR ONLY', 'font': SANS, 'size': 18, 'color': MUTED,
          'bold': True, 'letter_spacing': 240}], align=PP_ALIGN.CENTER)
    # £97
    add_paragraphs(s, Inches(2.5), Inches(2.6), Inches(8.3), Inches(2), [
        {'runs': [
            {'text': '£', 'font': SERIF, 'size': 60, 'color': CORAL, 'italic': True},
            {'text': '97', 'font': SERIF, 'size': 140, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.0},
    ])
    add_runs(s, Inches(2.5), Inches(5.0), Inches(8.3), Inches(0.5),
        [{'text': 'one-time payment  ·  12 months of WLA',
          'font': SANS, 'size': 17, 'color': MUTED, 'bold': True, 'letter_spacing': 200}],
        align=PP_ALIGN.CENTER)
    add_runs(s, Inches(2.5), Inches(5.5), Inches(8.3), Inches(0.5),
        [{'text': 'That’s less than £1.90 per week.',
          'font': SERIF, 'size': 22, 'color': CORAL, 'italic': True}], align=PP_ALIGN.CENTER)
    add_runs(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
        [{'text': 'The LOWEST price WLA will ever be offered at.',
          'font': SERIF, 'size': 18, 'color': NAVY, 'italic': True}], align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 46. NOT ANOTHER DIET (CALLOUT) ----------------
def slide_not_another_diet():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_runs(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.4),
        [{'text': 'REMEMBER', 'font': SANS, 'size': 18,
          'color': CORAL, 'bold': True, 'letter_spacing': 240}], align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'This isn’t about starting ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
            {'text': 'another diet.', 'font': SERIF, 'size': 44, 'color': CORAL, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    add_runs(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.6),
        [{'text': 'It’s about continuing the habits, support and momentum you’ve already worked so hard to build.',
          'font': SANS, 'size': 18, 'color': BODY}],
        align=PP_ALIGN.CENTER)
    add_runs(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6),
        [{'text': 'The women who maintain their results long term are usually the women who continue living the WLA lifestyle.',
          'font': SERIF, 'size': 19, 'color': NAVY, 'italic': True}],
        align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 47. IMAGINE YOUR SUMMER ----------------
def slide_imagine_summer():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    slide_eyebrow_headline(s, 'Imagine this being your summer',
        [{'text': 'Reaching September still ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
         {'text': 'building momentum.', 'font': SERIF, 'size': 40, 'color': CORAL, 'italic': True}],
        top=Inches(0.85))
    bullets = [
        'Feeling lighter',
        'More confident',
        'Back in control',
        'Reducing cravings',
        'Building routines that actually stick',
        'And reaching September STILL building momentum instead of starting over again',
    ]
    for i, b in enumerate(bullets):
        col = i % 2; row = i // 2
        x = Inches(0.9) if col == 0 else Inches(7.0)
        y = Inches(3.4) + Inches(0.62 * row)
        coral_bullet(s, x, y, b, size=16, color=NAVY)
    add_footer(s)

# ---------------- 48. NOT BECAUSE PERFECT ----------------
def slide_not_perfect():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM_2)
    add_paragraphs(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.5), [
        {'runs': [
            {'text': 'Not because you were ', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
            {'text': 'perfect', 'font': SERIF, 'size': 36, 'color': CORAL, 'italic': True},
            {'text': '.', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2, 'space_after': 18},
        {'runs': [
            {'text': 'But because you finally had:', 'font': SERIF, 'size': 24, 'color': NAVY, 'italic': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.2},
    ])
    pillars = ['Support', 'Accountability', 'Structure', 'The lifestyle']
    col_w_in = 2.65; h_in = 1.0; gap_in = 0.15
    total_w_in = col_w_in * 4 + gap_in * 3
    start_x_in = (13.333 - total_w_in) / 2
    y = Inches(5.0)
    for i, p in enumerate(pillars):
        x = Inches(start_x_in + (col_w_in + gap_in) * i)
        col_w = Inches(col_w_in); h = Inches(h_in)
        card = add_round_rect(s, x, y, col_w, h, PAPER, radius_pct=0.2)
        card.line.color.rgb = CORAL
        add_runs(s, x, y, col_w, h,
            [{'text': p, 'font': SERIF, 'size': 18, 'color': NAVY, 'bold': True}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s)

# ---------------- 49. YOUR CHANCE ----------------
def slide_your_chance():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, CREAM)
    add_pill(s, Inches(3.7), Inches(1.0), Inches(5.9), Inches(0.45),
             '✨  YOUR CHANCE TO BE PART OF THE NEXT CHAPTER  ✨',
             fill=NAVY, color=CREAM_WARM, size=15, letter_spacing=240)
    add_paragraphs(s, Inches(0.5), Inches(2.3), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'From the ', 'font': SERIF, 'size': 50, 'color': NAVY, 'bold': True},
            {'text': 'very beginning', 'font': SERIF, 'size': 50, 'color': CORAL, 'italic': True},
            {'text': '.', 'font': SERIF, 'size': 50, 'color': NAVY, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_runs(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.6),
        [{'text': 'Not just losing weight. Building a lifestyle that helps you KEEP the results too.',
          'font': SERIF, 'size': 21, 'color': NAVY, 'italic': True}], align=PP_ALIGN.CENTER)
    add_runs(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.5),
        [{'text': 'I truly believe this next phase of WLA is going to change everything for so many women.',
          'font': SANS, 'size': 17, 'color': BODY}], align=PP_ALIGN.CENTER)
    add_footer(s)

# ---------------- 50. CLOSE ----------------
def slide_close():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_pill(s, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.45),
             'FOUNDING MEMBERS · LIMITED 1,000 PLACES',
             fill=CORAL, color=WHITE, size=15, letter_spacing=260)
    add_paragraphs(s, Inches(0.5), Inches(2.5), Inches(12.4), Inches(3.5), [
        {'runs': [
            {'text': 'And I cannot ', 'font': SERIF, 'size': 56, 'color': WHITE, 'bold': True},
            {'text': 'wait', 'font': SERIF, 'size': 56, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
        {'runs': [
            {'text': 'for you to be part of it.', 'font': SERIF, 'size': 56, 'color': WHITE, 'bold': True},
        ], 'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    # CTA pill
    cta = add_round_rect(s, Inches(4.4), Inches(5.6), Inches(4.5), Inches(0.75),
                         CORAL, radius_pct=0.5)
    add_runs(s, Inches(4.4), Inches(5.6), Inches(4.5), Inches(0.75),
        [{'text': 'JOIN AS A FOUNDING MEMBER  →',
          'font': SANS, 'size': 18, 'color': WHITE, 'bold': True, 'letter_spacing': 240}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_runs(s, Inches(0.5), Inches(6.6), Inches(12.4), Inches(0.4),
        [{'text': 'sales.thewlacademy.com/wla-app-1yr',
          'font': SANS, 'size': 16, 'color': CREAM_WARM, 'letter_spacing': 150}],
        align=PP_ALIGN.CENTER)
    add_footer(s, dark=True)

# ============================================================
# BUILD
# ============================================================
slide_cover()             # 1
slide_opening()           # 2
slide_11_years()          # 3
slide_hardest_part()      # 4
slide_cycle()             # 5
slide_keep_it_off()       # 6
slide_pillars()           # 7
slide_evolution()         # 8
slide_shift()             # 9
slide_quote_lifestyle()   # 10
slide_i_live_it()         # 11
slide_not_perfectly()     # 12
slide_i_still()           # 13
slide_never_diet()        # 14
slide_taught_lived()      # 15
slide_passionate()        # 16
slide_welcome()           # 17
slide_about()             # 18
slide_huge()              # 19
slide_reveal_app()        # 20
slide_not_calorie()       # 21
slide_inside_overview()   # 22
slide_phase1_intro()      # 23
total_p1 = len(PHASE1_FEATURES)
for i, f in enumerate(PHASE1_FEATURES, start=1):  # 24..33
    feature_slide(i, total_p1, f['title'], f['lead'], f['bullets'], f.get('footer'))
slide_plus_beginning()    # 34
slide_phase2_intro()      # 35
for i, f in enumerate(PHASE2_FEATURES, start=1):  # 36..39
    phase2_slide(i, f['title'], f['lead'], f['bullets'])
slide_next_evolution()    # 40
slide_founding_open()     # 41
slide_gone_forever()      # 42
slide_only_opportunity()  # 43
slide_what_receive()      # 44
slide_pricing()           # 45
slide_not_another_diet()  # 46
slide_imagine_summer()    # 47
slide_not_perfect()       # 48
slide_your_chance()       # 49
slide_close()             # 50

out = '/home/user/21-day-reset/wla-summer-challenge-launch.pptx'
prs.save(out)
print('SAVED', out, 'slides=', len(prs.slides.__iter__.__self__._sldIdLst))
