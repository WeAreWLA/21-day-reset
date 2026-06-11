"""
Build the WLA VIP 1-2-1 Personal Coaching deck.
16:9 widescreen, brand-matched to the /12-wk-summer landing page:
  - cream paper background, navy serif headlines, coral italic accents
  - Libre Baskerville (serif) + Alegreya Sans (body)
  - navy callout slides, coral pill eyebrows, value/price cards
Same/similar copy as the existing VIP deck, adapted: flat pricing
(£597 / £1197), no 'relaunch' framing, evergreen + placeholder dates.

Output: wla-vip-121-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

# ---------- Brand ----------
NAVY        = RGBColor(0x0E, 0x27, 0x46)
NAVY_DEEP   = RGBColor(0x09, 0x1A, 0x30)
CORAL       = RGBColor(0xD9, 0x77, 0x57)
CORAL_SOFT  = RGBColor(0xE5, 0x91, 0x7A)
CREAM       = RGBColor(0xF9, 0xF7, 0xF4)
CREAM_2     = RGBColor(0xE9, 0xDF, 0xD3)
PAPER       = RGBColor(0xFD, 0xFB, 0xF8)
INK         = RGBColor(0x1A, 0x24, 0x38)
MUTED       = RGBColor(0x5A, 0x64, 0x78)
BODY        = RGBColor(0x33, 0x33, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_WARM  = RGBColor(0xF4, 0xEB, 0xDC)
RULE        = RGBColor(0xD9, 0xCD, 0xB6)

SERIF = 'Libre Baskerville'
SANS  = 'Alegreya Sans'

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

# ---------- Helpers (from build_deck.py design system) ----------
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    return bg

def add_runs(slide, left, top, width, height, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = line
    for run in runs:
        r = p.add_run(); r.text = run['text']
        r.font.name = run.get('font', SANS); r.font.size = Pt(run.get('size', 18))
        r.font.bold = run.get('bold', False); r.font.italic = run.get('italic', False)
        r.font.color.rgb = run.get('color', BODY)
        if 'letter_spacing' in run:
            r._r.get_or_add_rPr().set('spc', str(run['letter_spacing']))
    return tb

def add_paragraphs(slide, left, top, width, height, paragraphs,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align); p.line_spacing = para.get('line', 1.25)
        if 'space_after' in para: p.space_after = Pt(para['space_after'])
        if 'space_before' in para: p.space_before = Pt(para['space_before'])
        for run in para['runs']:
            r = p.add_run(); r.text = run['text']
            r.font.name = run.get('font', SANS); r.font.size = Pt(run.get('size', 18))
            r.font.bold = run.get('bold', False); r.font.italic = run.get('italic', False)
            r.font.color.rgb = run.get('color', BODY)
            if 'letter_spacing' in run:
                r._r.get_or_add_rPr().set('spc', str(run['letter_spacing']))
    return tb

def add_rect(slide, left, top, width, height, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None: s.line.fill.background()
    else: s.line.color.rgb = line_color
    s.shadow.inherit = False
    return s

def add_round_rect(slide, left, top, width, height, color, radius_pct=0.06,
                   line_color=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = radius_pct
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None: s.line.fill.background()
    else:
        s.line.color.rgb = line_color; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_pill(slide, left, top, width, height, text, fill=CORAL, color=WHITE, size=11,
             font=SANS, bold=True, letter_spacing=180):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = fill
    pill.line.fill.background(); pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    r._r.get_or_add_rPr().set('spc', str(letter_spacing))
    return pill

def coral_bullet(slide, left, top, text, size=18, font=SERIF, color=NAVY,
                 italic=False, bold=False, width=Inches(11.0), dot=CORAL):
    dot_d = Emu(Inches(0.13))
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Emu(int(Pt(size).emu * 0.45)), dot_d, dot_d)
    d.fill.solid(); d.fill.fore_color.rgb = dot
    d.line.fill.background(); d.shadow.inherit = False
    add_runs(slide, left + Inches(0.34), top, width, Inches(0.6),
        [{'text': text, 'font': font, 'size': size, 'color': color, 'italic': italic, 'bold': bold}])

def check_bullet(slide, left, top, label, body, width=Inches(5.2)):
    """Coral check + bold label + body, used for 'what's included'."""
    chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Emu(int(Pt(14).emu*0.2)),
                                 Emu(Inches(0.26)), Emu(Inches(0.26)))
    chk.fill.solid(); chk.fill.fore_color.rgb = CORAL; chk.line.fill.background()
    chk.shadow.inherit = False
    tf = chk.text_frame; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    rr = pr.add_run(); rr.text='✓'; rr.font.name=SANS; rr.font.size=Pt(11)
    rr.font.bold=True; rr.font.color.rgb=WHITE
    add_paragraphs(slide, left + Inches(0.44), top - Inches(0.02), width, Inches(1.1), [
        {'runs':[{'text':label,'font':SERIF,'size':18,'color':NAVY,'bold':True}], 'line':1.15,'space_after':2},
        {'runs':[{'text':body,'font':SANS,'size':15,'color':MUTED}], 'line':1.22},
    ])

def add_picture_cover(slide, path, left, top, width, height):
    iw, ih = Image.open(path).size
    target = width / height
    ar = iw / ih
    pic = slide.shapes.add_picture(path, left, top, width, height)
    if ar > target:
        crop = (1 - target / ar) / 2.0
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - ar / target) / 2.0
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic

def add_footer(slide, dark=False):
    return  # footer text removed per request
    color = CREAM_WARM if dark else MUTED
    add_runs(slide, Inches(0.7), Inches(6.98), Inches(8), Inches(0.4),
        [{'text': 'The Weight Loss Academy', 'font': SERIF, 'size': 12.5, 'color': color, 'italic': True}])
    add_runs(slide, Inches(7.6), Inches(6.98), Inches(5.0), Inches(0.4),
        [{'text': 'VIP 1-2-1 PERSONAL COACHING', 'font': SANS, 'size': 11,
          'color': color, 'bold': True, 'letter_spacing': 200}], align=PP_ALIGN.RIGHT)

def eyebrow(slide, text, x=Inches(0.9), y=Inches(0.82), color=CORAL):
    add_runs(slide, x, y, Inches(11.5), Inches(0.5),
        [{'text': text.upper(), 'font': SANS, 'size': 18, 'color': color,
          'bold': True, 'letter_spacing': 240}])

def heading(slide, runs, x=Inches(0.9), y=Inches(1.45), w=Inches(11.5), line=1.12):
    add_paragraphs(slide, x, y, w, Inches(2.0), [{'runs': runs, 'line': line, 'align': PP_ALIGN.LEFT}])

def coral_rule(slide, x=Inches(0.9), y=Inches(1.32), w=Inches(1.6)):
    add_rect(slide, x, y, w, Inches(0.055), CORAL)

# ============================================================
# SLIDES
# ============================================================

def s_cover():
    s = add_slide(); add_bg(s, CREAM)
    # right photo panel
    px = Inches(8.55); pw = Inches(13.333) - px
    add_picture_cover(s, os.path.join(ASSETS, 'anna-portrait.jpg'), px, 0, pw, SH)
    add_rect(s, px - Inches(0.06), 0, Inches(0.06), SH, CORAL)  # coral seam
    # left text
    add_pill(s, Inches(0.9), Inches(0.95), Inches(3.7), Inches(0.55),
             'VIP 1-2-1 COACHING', fill=NAVY, color=CREAM_WARM, size=14, letter_spacing=240)
    add_paragraphs(s, Inches(0.88), Inches(1.95), Inches(7.4), Inches(2.8), [
        {'runs': [{'text': 'Personal coaching', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True}], 'line': 1.12},
        {'runs': [{'text': 'that finally gets', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True}], 'line': 1.12},
        {'runs': [{'text': 'you ', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
                  {'text': 'results', 'font': SERIF, 'size': 42, 'color': CORAL, 'bold': True, 'italic': True},
                  {'text': '.', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True}], 'line': 1.12},
    ])
    add_runs(s, Inches(0.9), Inches(4.75), Inches(7.0), Inches(0.9),
        [{'text': 'Lose up to 3 stone — without extreme diets or exhausting workouts.',
          'font': SERIF, 'size': 20, 'color': MUTED, 'italic': True}], line=1.35)
    add_runs(s, Inches(0.9), Inches(6.35), Inches(7.0), Inches(0.5),
        [{'text': 'WITH ANNA WALLACE  ·  FOUNDER OF WLA  ·  REGISTERED ASSOCIATE NUTRITIONIST',
          'font': SANS, 'size': 11.5, 'color': NAVY, 'bold': True, 'letter_spacing': 150}])

def s_opening():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'A few VIP spots have opened'); coral_rule(s)
    heading(s, [
        {'text': 'I only coach a handful of women ', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
        {'text': '1-2-1', 'font': SERIF, 'size': 36, 'color': CORAL, 'bold': True, 'italic': True},
        {'text': ' at a time — and a few spots have ', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
        {'text': 'just opened', 'font': SERIF, 'size': 36, 'color': CORAL, 'bold': True, 'italic': True},
        {'text': '.', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
    ], y=Inches(1.55))
    add_paragraphs(s, Inches(0.9), Inches(3.7), Inches(11.0), Inches(3.0), [
        {'runs': [{'text': 'If you’ve been thinking about going all-in with personal coaching, ',
                   'font': SANS, 'size': 21, 'color': BODY},
                  {'text': 'this is your moment to grab one.',
                   'font': SANS, 'size': 21, 'color': NAVY, 'bold': True}], 'line': 1.5, 'space_after': 12},
        {'runs': [{'text': 'My full attention, daily support, and a plan built entirely around you.',
                   'font': SERIF, 'size': 22, 'color': NAVY, 'italic': True}], 'line': 1.4},
    ])
    add_footer(s)

def s_why_back():
    s = add_slide(); add_bg(s, NAVY)
    add_pill(s, Inches(4.82), Inches(1.45), Inches(3.7), Inches(0.5),
             'WHY VIP WORKS', fill=CORAL, color=WHITE, size=18, letter_spacing=200)
    add_paragraphs(s, Inches(1.2), Inches(2.55), Inches(10.9), Inches(3.0), [
        {'runs': [{'text': 'The difference isn’t the ', 'font': SERIF, 'size': 46, 'color': WHITE, 'italic': True},
                  {'text': 'plan', 'font': SERIF, 'size': 46, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
                  {'text': '.', 'font': SERIF, 'size': 46, 'color': WHITE, 'italic': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.15, 'space_after': 14},
        {'runs': [{'text': 'It’s having ', 'font': SERIF, 'size': 46, 'color': WHITE, 'italic': True},
                  {'text': 'me in your corner', 'font': SERIF, 'size': 46, 'color': CORAL_SOFT, 'italic': True, 'bold': True},
                  {'text': ' — every single day.', 'font': SERIF, 'size': 46, 'color': WHITE, 'italic': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.15},
    ])
    add_runs(s, Inches(1.2), Inches(5.55), Inches(10.9), Inches(0.7),
        [{'text': 'Personalised strategy. Daily support. Real accountability.',
          'font': SANS, 'size': 20, 'color': CREAM_WARM}], align=PP_ALIGN.CENTER)
    add_footer(s, dark=True)

def s_what_vip():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'What VIP helps you do'); coral_rule(s)
    heading(s, [{'text': 'Lose the weight — and ', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True},
                {'text': 'feel like you again', 'font': SERIF, 'size': 42, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': '.', 'font': SERIF, 'size': 42, 'color': NAVY, 'bold': True}])
    items = ['Lose up to 1–3 stone', 'Build real, lasting confidence',
             'Calm your sugar cravings', 'Skyrocket your energy',
             'Look and feel amazing by summer', 'No extreme diets. No exhausting workouts.']
    y = Inches(3.5)
    for i, it in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.95) if col == 0 else Inches(6.85)
        coral_bullet(s, x, y + Inches(0.78 * row), it, size=20, color=NAVY, width=Inches(5.6))
    add_footer(s)

def s_difference():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Why 1-2-1 changes everything'); coral_rule(s)
    heading(s, [{'text': 'Nothing transforms results like ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
                {'text': '1-2-1 coaching', 'font': SERIF, 'size': 40, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': '.', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True}])
    add_paragraphs(s, Inches(0.9), Inches(3.35), Inches(11.2), Inches(3.2), [
        {'runs': [{'text': 'I’ll guide you through the exact simple nutrition formula I used to lose over 2 stone and keep it off.',
                   'font': SANS, 'size': 21, 'color': BODY}], 'line': 1.5, 'space_after': 12},
        {'runs': [{'text': 'As a degree-qualified Registered Associate Nutritionist with 11+ years’ experience and 50,000+ clients supported — ',
                   'font': SANS, 'size': 21, 'color': BODY},
                  {'text': 'nothing comes close to personal 1-2-1 coaching.',
                   'font': SERIF, 'size': 20, 'color': NAVY, 'italic': True}], 'line': 1.5},
    ])
    add_footer(s)

def s_steps_intro():
    s = add_slide(); add_bg(s, NAVY)
    add_pill(s, Inches(4.92), Inches(1.4), Inches(3.5), Inches(0.5),
             'THE VIP METHOD', fill=CORAL, color=WHITE, size=18, letter_spacing=200)
    add_paragraphs(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.0), [
        {'runs': [{'text': 'Three things that make VIP ', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True},
                  {'text': 'different', 'font': SERIF, 'size': 46, 'color': CORAL_SOFT, 'bold': True, 'italic': True},
                  {'text': '.', 'font': SERIF, 'size': 46, 'color': WHITE, 'bold': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.1}])
    labels = [('01', 'Personalised Strategy'), ('02', 'High Accountability'), ('03', 'Results Tracking')]
    cw = Inches(3.5); gap = Inches(0.45); total = cw * 3 + gap * 2
    x0 = (SW - total) / 2
    for i, (n, t) in enumerate(labels):
        x = x0 + (cw + gap) * i
        add_round_rect(s, x, Inches(4.0), cw, Inches(2.0), NAVY_DEEP, radius_pct=0.08, line_color=CORAL, line_w=1.0)
        add_runs(s, x, Inches(4.25), cw, Inches(0.9),
            [{'text': n, 'font': SERIF, 'size': 44, 'color': CORAL_SOFT, 'bold': True}], align=PP_ALIGN.CENTER)
        add_runs(s, x + Inches(0.2), Inches(5.25), cw - Inches(0.4), Inches(0.8),
            [{'text': t, 'font': SERIF, 'size': 20, 'color': WHITE, 'italic': True}], align=PP_ALIGN.CENTER, line=1.15)
    add_footer(s, dark=True)

def s_step(num, title, lead, paras):
    s = add_slide(); add_bg(s, CREAM)
    # watermark number
    add_runs(s, Inches(9.7), Inches(0.5), Inches(3.4), Inches(2.4),
        [{'text': f'{num:02d}', 'font': SERIF, 'size': 150, 'color': CREAM_2, 'bold': True}], align=PP_ALIGN.RIGHT)
    eyebrow(s, f'Step {num} of 3'); coral_rule(s)
    heading(s, [{'text': title, 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True}],
            w=Inches(9.2), line=1.12)
    add_runs(s, Inches(0.9), Inches(3.3), Inches(9.6), Inches(0.7),
        [{'text': lead, 'font': SERIF, 'size': 21, 'color': CORAL, 'italic': True}], line=1.3)
    add_paragraphs(s, Inches(0.9), Inches(4.25), Inches(11.2), Inches(2.4),
        [{'runs': [{'text': p, 'font': SANS, 'size': 20, 'color': BODY}], 'line': 1.5, 'space_after': 8} for p in paras])
    add_footer(s)

def s_stanford():
    s = add_slide(); add_bg(s, NAVY)
    eyebrow(s, 'The research', color=CORAL_SOFT);
    add_rect(s, Inches(0.9), Inches(1.32), Inches(1.6), Inches(0.055), CORAL)
    add_paragraphs(s, Inches(0.9), Inches(1.55), Inches(11.4), Inches(2.6), [
        {'runs': [{'text': 'Human 1-2-1 coaching ', 'font': SERIF, 'size': 42, 'color': WHITE, 'bold': True},
                  {'text': 'outperforms', 'font': SERIF, 'size': 42, 'color': CORAL_SOFT, 'bold': True, 'italic': True},
                  {'text': ' digital-only or self-guided approaches.', 'font': SERIF, 'size': 42, 'color': WHITE, 'bold': True}],
         'line': 1.14}])
    add_paragraphs(s, Inches(0.9), Inches(4.4), Inches(11.4), Inches(2.2), [
        {'runs': [{'text': 'A Stanford University study found people with access to a human coach lost significantly more weight than those using digital tools alone — even when the information was identical.',
                   'font': SANS, 'size': 21, 'color': CREAM_WARM}], 'line': 1.5, 'space_after': 12},
        {'runs': [{'text': 'It’s not the plan. It’s the support, accountability and feedback.',
                   'font': SERIF, 'size': 21, 'color': WHITE, 'italic': True}], 'line': 1.4},
    ])
    add_footer(s, dark=True)

def s_stay_stuck():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Why most women stay stuck'); coral_rule(s)
    heading(s, [{'text': 'You already know what to do. ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
                {'text': 'Following through is the hard part.', 'font': SERIF, 'size': 38, 'color': CORAL, 'bold': True, 'italic': True}])
    add_paragraphs(s, Inches(0.9), Inches(3.5), Inches(11.2), Inches(3.0), [
        {'runs': [{'text': 'Most women are stuck in the cycle of starting… stopping… and never quite following through. A good week, then life gets busy, motivation drops, and it’s back to square one.',
                   'font': SANS, 'size': 21, 'color': BODY}], 'line': 1.5, 'space_after': 12},
        {'runs': [{'text': 'Not because they don’t care — because they’re trying to do it alone. That’s exactly why VIP works.',
                   'font': SERIF, 'size': 20, 'color': NAVY, 'italic': True}], 'line': 1.45},
    ])
    add_footer(s)

def s_imagine():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Imagine, by summer'); coral_rule(s)
    heading(s, [{'text': 'Picture this version of ', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True},
                {'text': 'you', 'font': SERIF, 'size': 40, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': '…', 'font': SERIF, 'size': 40, 'color': NAVY, 'bold': True}])
    items = ['Putting on your favourite clothes and loving how you look',
             'Waking up lighter, energised and comfortable',
             'No more constant cravings or evening snacking',
             'Feeling proud of the woman in the mirror',
             'Enjoying holidays and BBQs without worrying about your weight',
             'Knowing exactly how to eat to lose weight and keep it off']
    y = Inches(3.35)
    for i, it in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.95) if col == 0 else Inches(6.85)
        coral_bullet(s, x, y + Inches(0.82 * row), it, size=19, font=SANS, color=BODY, width=Inches(5.7))
    add_footer(s)

def s_struggling():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Sound familiar?'); coral_rule(s)
    heading(s, [{'text': 'Are you ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
                {'text': 'struggling with', 'font': SERIF, 'size': 38, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': '…?', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True}], y=Inches(1.4))
    cards = [
        ('Time', 'No time to figure out your nutrition? Get a personalised strategy done for you.'),
        ('Results', 'Sick of not seeing results and want monthly progress you can feel?'),
        ('Consistency', 'Need a coach so you actually stay consistent instead of quitting?'),
        ('Energy', 'Want a setup that gives you energy — off the sugar-and-caffeine cycle?'),
        ('Snacking', 'Need support to stop the evening snacking that’s just habit?'),
        ('Weekends', 'Want weekend handholding so you don’t undo your hard work?'),
    ]
    cw = Inches(3.7); ch = Inches(1.95); gx = Inches(0.3); gy = Inches(0.3)
    x0 = Inches(0.9); y0 = Inches(2.55)
    for i, (lab, body) in enumerate(cards):
        col = i % 3; row = i // 3
        x = x0 + (cw + gx) * col; y = y0 + (ch + gy) * row
        add_round_rect(s, x, y, cw, ch, WHITE, radius_pct=0.07, line_color=CREAM_2, line_w=1.0)
        add_runs(s, x + Inches(0.28), y + Inches(0.22), cw - Inches(0.5), Inches(0.4),
            [{'text': lab.upper(), 'font': SANS, 'size': 14, 'color': CORAL, 'bold': True, 'letter_spacing': 150}])
        add_paragraphs(s, x + Inches(0.28), y + Inches(0.66), cw - Inches(0.56), ch - Inches(0.8),
            [{'runs': [{'text': body, 'font': SANS, 'size': 15, 'color': BODY}], 'line': 1.3}])
    add_footer(s)

def s_jill():
    s = add_slide(); add_bg(s, CREAM)
    # photo left
    pw = Inches(5.4)
    add_picture_cover(s, os.path.join(ASSETS, 'jill-before-after.jpg'), 0, 0, pw, SH)
    add_rect(s, pw, 0, Inches(0.06), SH, CORAL)
    # text right
    tx = Inches(6.0)
    add_pill(s, tx, Inches(1.2), Inches(2.7), Inches(0.5), 'VIP SUCCESS STORY',
             fill=NAVY, color=CREAM_WARM, size=12, letter_spacing=200)
    add_paragraphs(s, tx, Inches(2.1), Inches(6.6), Inches(2.6), [
        {'runs': [{'text': 'Jill lost ', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True},
                  {'text': '4 stone', 'font': SERIF, 'size': 44, 'color': CORAL, 'bold': True, 'italic': True}], 'line': 1.1},
        {'runs': [{'text': 'in 6 months in VIP.', 'font': SERIF, 'size': 44, 'color': NAVY, 'bold': True}], 'line': 1.1},
    ])
    add_paragraphs(s, tx, Inches(4.5), Inches(6.5), Inches(2.0), [
        {'runs': [{'text': '“VIP has changed my life. I did WLA resets and group for 2 years and didn’t achieve half of what I have in VIP. It’s just so different.”',
                   'font': SERIF, 'size': 20, 'color': MUTED, 'italic': True}], 'line': 1.45}])
    add_footer(s)

def s_results_montage():
    s = add_slide(); add_bg(s, NAVY)
    eyebrow(s, 'Real women, real results', color=CORAL_SOFT)
    add_rect(s, Inches(0.9), Inches(1.32), Inches(1.6), Inches(0.055), CORAL)
    add_paragraphs(s, Inches(0.9), Inches(1.5), Inches(11.4), Inches(1.2), [
        {'runs': [{'text': 'The transformations speak for ', 'font': SERIF, 'size': 36, 'color': WHITE, 'bold': True},
                  {'text': 'themselves', 'font': SERIF, 'size': 36, 'color': CORAL_SOFT, 'bold': True, 'italic': True},
                  {'text': '.', 'font': SERIF, 'size': 36, 'color': WHITE, 'bold': True}], 'line': 1.1}])
    imgs = ['ruth-before-after.jpg', 'laura-before-after.png', 'vicky-before-after.png', 'barbara-before-after.png']
    iw = Inches(2.78); gap = Inches(0.25); total = iw * 4 + gap * 3
    x0 = (SW - total) / 2; y = Inches(3.0)
    for i, im in enumerate(imgs):
        x = x0 + (iw + gap) * i
        add_picture_cover(s, os.path.join(ASSETS, im), x, y, iw, iw)
    add_runs(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.5),
        [{'text': 'From 1 to 4 stone down — in 1-2-1 VIP coaching with Anna.',
          'font': SERIF, 'size': 19, 'color': CREAM_WARM, 'italic': True}], align=PP_ALIGN.CENTER)
    add_footer(s, dark=True)

def s_included():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Everything you get'); coral_rule(s)
    heading(s, [{'text': 'What’s ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
                {'text': 'included', 'font': SERIF, 'size': 38, 'color': CORAL, 'bold': True, 'italic': True}], y=Inches(1.42))
    items = [
        ('Monthly 1-2-1 Coaching Sessions', 'A dedicated session each month to review progress and refine your plan.'),
        ('Daily WhatsApp / Messenger Support', 'Your coach in your pocket for guidance and encouragement every day.'),
        ('Weekly Food Diary Review', 'We review your nutrition each week and adjust what’s working.'),
        ('Weekly Nutrition Strategy', 'Your plan evolves with your progress, lifestyle and goals.'),
        ('Laser-Focused Weekly Actions', 'Clear priorities each week so you always know your next step.'),
        ('Monthly Progress Tracking & Review', 'Structured check-ins to keep you accountable and moving forward.'),
        ('Shift & Sustain Members Area', 'Full access to meal guides, recipes and the WLA framework.'),
        ('A Strategy Built Around You', 'No guesswork, no extremes — just a clear plan that fits your life.'),
    ]
    x0 = Inches(0.9); y0 = Inches(2.7); rowh = Inches(1.05)
    for i, (lab, body) in enumerate(items):
        col = i % 2; row = i // 2
        x = x0 + (Inches(6.0)) * col
        check_bullet(s, x, y0 + rowh * row, lab, body, width=Inches(5.3))
    add_footer(s)

def s_investment():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Investment vs reality'); coral_rule(s)
    heading(s, [{'text': 'This level of support usually costs ', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True},
                {'text': 'thousands', 'font': SERIF, 'size': 36, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': '.', 'font': SERIF, 'size': 36, 'color': NAVY, 'bold': True}])
    add_paragraphs(s, Inches(0.9), Inches(3.4), Inches(11.2), Inches(3.0), [
        {'runs': [{'text': 'Most degree-qualified Registered Associate Nutritionists charge £300–£700 for consultations and ongoing monthly support.',
                   'font': SANS, 'size': 21, 'color': BODY}], 'line': 1.5, 'space_after': 10},
        {'runs': [{'text': 'True 1-2-1 coaching — daily support, weekly reviews, a personalised strategy — typically runs £2,000–£5,000+ (take a look around Harley Street).',
                   'font': SANS, 'size': 21, 'color': BODY}], 'line': 1.5, 'space_after': 10},
        {'runs': [{'text': 'For this round, I’ve made it far more accessible.',
                   'font': SERIF, 'size': 21, 'color': NAVY, 'italic': True}], 'line': 1.4},
    ])
    add_footer(s)

def s_pricing():
    s = add_slide(); add_bg(s, CREAM)
    add_runs(s, Inches(0.9), Inches(0.62), Inches(11.5), Inches(0.5),
        [{'text': 'CHOOSE YOUR VIP JOURNEY', 'font': SANS, 'size': 18, 'color': CORAL,
          'bold': True, 'letter_spacing': 240}], align=PP_ALIGN.CENTER)
    add_runs(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(0.8),
        [{'text': 'Personal 1-2-1 coaching, built around you.', 'font': SERIF, 'size': 30,
          'color': NAVY, 'bold': True, 'italic': False}], align=PP_ALIGN.CENTER)

    cw = Inches(5.3); ch = Inches(3.55); gap = Inches(0.6)
    x0 = (SW - (cw * 2 + gap)) / 2; y = Inches(2.5)

    # 3 months card
    add_round_rect(s, x0, y, cw, ch, WHITE, radius_pct=0.05, line_color=CREAM_2, line_w=1.25)
    add_runs(s, x0, y + Inches(0.35), cw, Inches(0.5),
        [{'text': 'VIP · 3 MONTHS', 'font': SANS, 'size': 14, 'color': NAVY, 'bold': True, 'letter_spacing': 200}], align=PP_ALIGN.CENTER)
    add_runs(s, x0, y + Inches(0.9), cw, Inches(1.2),
        [{'text': '£597', 'font': SERIF, 'size': 64, 'color': CORAL, 'bold': True}], align=PP_ALIGN.CENTER)
    add_runs(s, x0 + Inches(0.5), y + Inches(2.15), cw - Inches(1.0), Inches(0.7),
        [{'text': '3 months of personalised 1-2-1 coaching', 'font': SANS, 'size': 17, 'color': BODY}], align=PP_ALIGN.CENTER, line=1.3)
    add_runs(s, x0 + Inches(0.5), y + Inches(2.95), cw - Inches(1.0), Inches(0.5),
        [{'text': 'Monthly plans from £199', 'font': SERIF, 'size': 17, 'color': NAVY, 'italic': True}], align=PP_ALIGN.CENTER)

    # 6 months card (highlighted)
    x1 = x0 + cw + gap
    add_round_rect(s, x1, y, cw, ch, NAVY, radius_pct=0.05)
    add_pill(s, x1 + cw/2 - Inches(1.1), y - Inches(0.28), Inches(2.2), Inches(0.5),
             'BEST VALUE', fill=CORAL, color=WHITE, size=12, letter_spacing=200)
    add_runs(s, x1, y + Inches(0.35), cw, Inches(0.5),
        [{'text': 'VIP · 6 MONTHS', 'font': SANS, 'size': 14, 'color': CREAM_WARM, 'bold': True, 'letter_spacing': 200}], align=PP_ALIGN.CENTER)
    add_runs(s, x1, y + Inches(0.9), cw, Inches(1.2),
        [{'text': '£1197', 'font': SERIF, 'size': 64, 'color': CORAL_SOFT, 'bold': True}], align=PP_ALIGN.CENTER)
    add_runs(s, x1 + Inches(0.5), y + Inches(2.15), cw - Inches(1.0), Inches(0.7),
        [{'text': 'The time to lose it AND keep it off for good', 'font': SANS, 'size': 17, 'color': CREAM_WARM}], align=PP_ALIGN.CENTER, line=1.3)
    add_runs(s, x1 + Inches(0.5), y + Inches(2.95), cw - Inches(1.0), Inches(0.5),
        [{'text': 'Most of my biggest transformations', 'font': SERIF, 'size': 17, 'color': WHITE, 'italic': True}], align=PP_ALIGN.CENTER)

    add_runs(s, Inches(0.6), Inches(6.18), Inches(12.13), Inches(1.1),
        [{'text': 'Pay with PayPal, Klarna or Clearpay  ·  Limited to 30 women  ·  Doors open [START DATE]',
          'font': SANS, 'size': 24, 'color': NAVY, 'bold': True, 'letter_spacing': 20}],
        align=PP_ALIGN.CENTER, line=1.25)

def s_who_for():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Is this you?'); coral_rule(s)
    heading(s, [{'text': 'Who VIP is ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
                {'text': 'for', 'font': SERIF, 'size': 38, 'color': CORAL, 'bold': True, 'italic': True}], y=Inches(1.42))
    items = ['You know you need support and accountability to stay consistent',
             'You’re tired of figuring it all out alone',
             'You take action when someone’s holding you to it',
             'You want a coach for when motivation dips or life gets busy',
             'You’re ready to follow through on your goals',
             'You want to lose up to 3 stone — sustainably, for good']
    y = Inches(2.7)
    for i, it in enumerate(items):
        coral_bullet(s, Inches(0.95), y + Inches(0.66 * i), it, size=19, font=SANS, color=BODY, width=Inches(11.3))
    add_footer(s)

def s_who_not():
    s = add_slide(); add_bg(s, CREAM)
    eyebrow(s, 'Honesty first'); coral_rule(s)
    heading(s, [{'text': 'Who VIP is ', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True},
                {'text': 'not', 'font': SERIF, 'size': 38, 'color': CORAL, 'bold': True, 'italic': True},
                {'text': ' for', 'font': SERIF, 'size': 38, 'color': NAVY, 'bold': True}], y=Inches(1.42))
    items = ['You’re looking for a quick fix or crash diet',
             'You’re not ready to take action or follow through',
             'You want results without changing any habits',
             'You’d rather do everything alone, without guidance',
             'You’re not open to a sustainable approach to nutrition',
             'You’re not ready to commit right now']
    y = Inches(2.7)
    for i, it in enumerate(items):
        coral_bullet(s, Inches(0.95), y + Inches(0.66 * i), it, size=19, font=SANS, color=MUTED, width=Inches(11.3), dot=CREAM_2)
    add_footer(s)

def s_two_options():
    s = add_slide(); add_bg(s, NAVY)
    add_pill(s, Inches(5.17), Inches(1.3), Inches(3.0), Inches(0.5),
             'TWO OPTIONS', fill=CORAL, color=WHITE, size=18, letter_spacing=200)
    add_paragraphs(s, Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.0), [
        {'runs': [{'text': 'You can keep doing this ', 'font': SERIF, 'size': 38, 'color': WHITE, 'bold': True},
                  {'text': 'alone', 'font': SERIF, 'size': 38, 'color': CORAL_SOFT, 'bold': True, 'italic': True},
                  {'text': ' — and be in the same place in 3–6 months.', 'font': SERIF, 'size': 38, 'color': WHITE, 'bold': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.18, 'space_after': 24},
        {'runs': [{'text': 'Or get the support, structure and accountability to ', 'font': SERIF, 'size': 38, 'color': WHITE, 'bold': True},
                  {'text': 'finally follow through', 'font': SERIF, 'size': 38, 'color': CORAL_SOFT, 'bold': True, 'italic': True},
                  {'text': '.', 'font': SERIF, 'size': 38, 'color': WHITE, 'bold': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.18}])
    add_footer(s, dark=True)

def s_cta():
    s = add_slide(); add_bg(s, NAVY)
    add_pill(s, Inches(5.52), Inches(0.95), Inches(2.3), Inches(0.5),
             'READY?', fill=CORAL, color=WHITE, size=18, letter_spacing=250)
    add_paragraphs(s, Inches(1.0), Inches(1.85), Inches(11.3), Inches(2.0), [
        {'runs': [{'text': 'Reply and say ', 'font': SERIF, 'size': 44, 'color': WHITE, 'bold': True},
                  {'text': '“I’m ready for VIP.”', 'font': SERIF, 'size': 44, 'color': CORAL_SOFT, 'bold': True, 'italic': True}],
         'align': PP_ALIGN.CENTER, 'line': 1.12}])
    add_paragraphs(s, Inches(1.4), Inches(3.7), Inches(10.5), Inches(1.6), [
        {'runs': [{'text': 'I’ll let you know if there are spaces left. Spaces are limited to 30 women — and previous VIP rounds have always filled.',
                   'font': SANS, 'size': 20, 'color': CREAM_WARM}], 'align': PP_ALIGN.CENTER, 'line': 1.45}])
    add_runs(s, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.6),
        [{'text': 'Email me: ', 'font': SANS, 'size': 20, 'color': WHITE, 'bold': True},
         {'text': 'annasnutrition@gmail.com', 'font': SANS, 'size': 20, 'color': CORAL_SOFT, 'bold': True}],
        align=PP_ALIGN.CENTER)
    add_runs(s, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.6),
        [{'text': 'Anna Wallace', 'font': SERIF, 'size': 19, 'color': WHITE, 'italic': True},
         {'text': '  ·  Founder of The Weight Loss Academy, Registered Associate Nutritionist',
          'font': SANS, 'size': 15, 'color': CREAM_WARM}], align=PP_ALIGN.CENTER)

# ---------- Build sequence ----------
s_cover()
s_opening()
s_why_back()
s_what_vip()
s_difference()
s_steps_intro()
s_step(1, 'Personalised Nutrition Strategy', 'A plan built entirely around you.',
       ['We build a simple nutrition strategy tailored to your lifestyle, goals, time and preferences.',
        'No guesswork. No extremes. No starting again every Monday — just a clear structure that reduces cravings and gets results sooner.'])
s_step(2, 'High Accountability & Daily Support', 'The part most women are missing.',
       ['When you’re accountable to someone else, you show up, follow through and take your goals seriously.',
        'With regular check-ins, coaching and daily support, you stay on track — even on the days you don’t feel like it.'])
s_step(3, 'Results Tracking', 'We measure what matters, every month.',
       ['You submit your results each month so we can track progress clearly and keep you moving forward.',
        'We review everything and make the right adjustments along the way, so your results keep progressing.'])
s_stanford()
s_stay_stuck()
s_imagine()
s_struggling()
s_jill()
s_results_montage()
s_included()
s_investment()
s_pricing()
s_who_for()
s_who_not()
s_two_options()
s_cta()

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'wla-vip-121-deck.pptx')
prs.save(out)
print('Saved', out, 'with', len(prs.slides._sldIdLst), 'slides')
